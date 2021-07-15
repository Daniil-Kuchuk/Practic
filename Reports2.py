from bs4 import BeautifulSoup as bs
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
import markdown


class Reports:
    def __init__(self, path):
        self.__path = path

    def create_presentation(self):
        markdown.markdownFromFile(input=self.__path, output='reports.html', extensions=['extra'])

        with open('reports.html') as fp:
            soup = bs(fp, 'html.parser')

        prs = Presentation()

        for el in soup.findAll('section'):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.width, title.height = Pt(720), Pt(60)
            title.text = el.h1.text
            p_list = el.findAll('p')
            if len(p_list):
                for p in p_list:
                    # for shape in slide.shapes:
                    #     if not shape.has_text_frame:
                    #         continue
                    #     text_frame = shape.text_frame
                    txt_box = slide.shapes.add_textbox(Pt(0), Pt(60), Pt(720), Pt(200))
            #         plc = slide.shapes.placeholders[1]
                    text_frame = txt_box.text_frame
                    text_frame.text = p.text
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            #         text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            #         # print(txt_box.width.pt)
            #         text_frame.word_wrap = True

        prs.save('reports.pptx')

