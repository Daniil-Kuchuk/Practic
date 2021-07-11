from pptx import Presentation
from pptx.util import Pt
import os.path
import sys


class Reports:
    def __init__(self, path='C:/', save_to='C:/'):
        self.__prefix = {}
        self.__path = path if os.path.exists(path) else None
        self.__save_to = save_to if os.path.exists(save_to) else None
        self.__images = {}

    def __stop_program(self):
        if self.__path is None:
            print("a;ldf")
            sys.exit(-1)

    def create_slide(self):
        self.__stop_program()

        self.__template()

        prs = Presentation()
        numbers_row = len(os.listdir(self.__path))
        numbers_col = len(os.listdir(os.chdir(self.__path + '\\' + os.listdir(self.__path)[0])))
        img_width = Pt((690 - 10 * (numbers_col - 1)) / numbers_col)
        img_height = Pt((450 - 10 * (numbers_row - 1)) / numbers_row)

        for key, value in self.__prefix.items():
            pos_y = Pt(90)
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title
            title.width, title.height = Pt(720), Pt(60)
            title.text = f'{value}'
            for grid in range(numbers_row):
                pos_x = Pt(30)
                txt_box = slide.shapes.add_textbox(Pt(0), pos_y, Pt(30), Pt(60))
                txt_box.text_frame._bodyPr.attrib.update({'vert': 'vert'})
                txt_box.text = f'сетка {grid+1}'
                for col in range(numbers_col):
                    q_txt_box = slide.shapes.add_textbox(pos_x, Pt(65), Pt(30), Pt(30))
                    q_txt_box.text = f'q{col+1}'
                    slide.shapes.add_picture(self.__images[f'{key}_grid{grid+1}_q{col+1}'],
                                             pos_x, pos_y,
                                             width=img_width,
                                             height=img_height)
                    pos_x += img_width + Pt(10)
                pos_y += img_height + Pt(10)

        prs.save(self.__save_to)

    def add_prefix(self, prefix):
        for (key, value) in prefix.items():
            self.__prefix[key] = value

    def __template(self):
        for i, grid in enumerate(os.listdir(self.__path)):
            for q in os.listdir(f'{self.__path}\\{grid}'):
                for img in os.listdir(f'{self.__path}\\{grid}\\{q}'):
                    self.__images[f'{img.split("_")[0]}_grid{i+1}_{q}'] = f'{self.__path}\\{grid}\\{q}\\{img}'
